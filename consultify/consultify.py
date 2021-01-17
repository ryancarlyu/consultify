import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.dml.line import LineFormat
from pptx.enum.dml import MSO_LINE
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.util import Pt

from pptx.oxml.xmlchemy import OxmlElement

def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element
    
def makeParaBulletPointed(para):
    """Bullets are set to Arial,
        actual text can be a different font"""
    pPr = para._p.get_or_add_pPr()
    ## Set marL and indent attributes
    pPr.set('marL','91440') #171450
    pPr.set('indent','171450') #171450
    ## Add buFont
    _ = SubElement(parent=pPr,
                   tagname="a:buFont",
                   typeface="Arial",
                   panose="020B0604020202020204",
                   pitchFamily="34",
                   charset="0"
                   )
    ## Add buChar
    _ = SubElement(parent=pPr,
                   tagname='a:buChar',
                   char="•")

def _set_cell_border(cell, border_color="000000", dash='solid', border_width='12700',transparent=False):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:noFill')
    for lines in ['a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        if transparent == True:
          solidFill = SubElement(ln, 'a:noFill')
        else:
          solidFill = SubElement(ln, 'a:solidFill')
          if border_color=='black':
            srgbClr = SubElement(solidFill, 'a:srgbClr', val="000000")
          elif border_color=='gray':
            srgbClr = SubElement(solidFill, 'a:srgbClr', val="D3D3D3")
          else:
            srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
          prstDash = SubElement(ln, 'a:prstDash', val=dash)
          round_ = SubElement(ln, 'a:round')
          headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
          tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

def make_deck():
  return Presentation()

def add_marvin_table_slide(presentation_name, df, cell_width=1.5, \
                           cell_height=0.5, slide_title="", cell_font_size=16, \
                           title_font_size=30):
  """
  Turns a pandas DataFrame into a consulting-style PowerPoint slide

  Parameters
  ----------  
  presentation_name : Presentation object
      The Presentation object created using the make_deck() function.
  
  df : pandas DataFrame

  cell_width : float, optional
      The width of each cell in inches (default is 1.5).
  
  cell_height : float, optional
      The height of each cell in inches (default is 0.5).

  slide_title : str, optional
      The text that serves as the title of the slide (default is no text).
  
  cell_font_size : int, optional
      The font size of text in each cell (default is 16).
  
  title_font_size : int, optional
      The font size of title text (default is 30).
  """
  slide = presentation_name.slides.add_slide(presentation_name.slide_layouts[5])

  ncols = df.shape[1]
  nrows = df.shape[0] + 1
  cellwidth = Inches(cell_width)
  cellheight = Inches(cell_height)
  bufferwidth = Inches(0.25)
  left = (presentation_name.slide_width-(ncols * cellwidth + (ncols-1) * bufferwidth))/2
  top = (presentation_name.slide_height-(nrows * cellheight))/2
  slide.shapes.title.text = slide_title
  slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
  slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(title_font_size)

  table = slide.shapes.add_table(rows=nrows, cols=ncols * 2 - 1, \
                                left=Inches(left/Inches(1)), \
                                top=Inches(top/Inches(1)), \
                                width=Inches((ncols * cellwidth + (ncols-1) * \
                                              bufferwidth)/Inches(1)), \
                                height=Inches((nrows * cellheight)/Inches(1)))

  for cols in range(ncols * 2 - 1):
    if cols % 2 == 1:
      table.table.columns[cols].width = bufferwidth
    else:
      table.table.columns[cols].width = cellwidth

  for rows in range(nrows):
    for cols in range(ncols * 2 - 1):
      if rows == 0:
        if cols % 2 == 1:
          _set_cell_border(table.table.cell(rows,cols),transparent=True)
          table.table.cell(rows,cols).fill.solid()
          table.table.cell(rows,cols).fill.background()
          continue
        _set_cell_border(table.table.cell(rows,cols))
        table.table.cell(rows,cols).text_frame.add_paragraph
        table.table.cell(rows,cols).fill.solid()
        table.table.cell(rows,cols).fill.background()
        table.table.cell(rows,cols).text_frame.paragraphs[0].text = df.columns[cols//2]
        table.table.cell(rows,cols).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x0,0x0,0x0)
        table.table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(cell_font_size)
        table.table.cell(rows,cols).text_frame.word_wrap = True
        table.table.cell(rows,cols).vertical_anchor = MSO_ANCHOR.BOTTOM
      elif rows == nrows - 1:
        _set_cell_border(table.table.cell(rows,cols),transparent=True)
        table.table.cell(rows,cols).text_frame.add_paragraph
        table.table.cell(rows,cols).fill.solid()
        table.table.cell(rows,cols).fill.background()
        if cols % 2 == 1:
          continue
        table.table.cell(rows,cols).text_frame.paragraphs[0].text = str(df.iloc[rows - 1,cols//2])
        makeParaBulletPointed(table.table.cell(rows,cols).text_frame.paragraphs[0])
        table.table.cell(rows,cols).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x0,0x0,0x0)
        table.table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(cell_font_size)
        table.table.cell(rows,cols).text_frame.word_wrap = True
        table.table.cell(rows,cols).vertical_anchor = MSO_ANCHOR.MIDDLE
      else:
        _set_cell_border(table.table.cell(rows,cols),border_color='gray',dash='dash')
        table.table.cell(rows,cols).text_frame.add_paragraph
        table.table.cell(rows,cols).fill.solid()
        table.table.cell(rows,cols).fill.background()
        if cols % 2 == 1:
          continue
        table.table.cell(rows,cols).text_frame.paragraphs[0].text = str(df.iloc[rows - 1,cols//2])
        makeParaBulletPointed(table.table.cell(rows,cols).text_frame.paragraphs[0])
        table.table.cell(rows,cols).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x0,0x0,0x0)
        table.table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(cell_font_size)
        table.table.cell(rows,cols).text_frame.word_wrap = True
        table.table.cell(rows,cols).vertical_anchor = MSO_ANCHOR.MIDDLE
        
  table.left = Inches((presentation_name.slide_width - table.width)/2/Inches(1))
  table.top = Inches((presentation_name.slide_height - table.height)/2/Inches(1))

def add_slide(presentation_name, slide_title = "", slide_text = "", image_filepath = None, \
              max_image = True, textbox_filled_space = 0.5, title_font_size = 30, \
              textbox_font_size = 16):
  """
  Creates a standard slide with a title, bulleted text box, and/or image.

  Parameters
  ----------
  presentation_name : Presentation object
      The Presentation object created using the make_deck() function.

  slide_title : str, optional
      The text that serves as the title of the slide (default is no text).
  
  slide_text : str, optional
      The bulleted text in the slide (default is no text).
  
  image_filepath : filepath, optional
      The filepath to the image file or the filename if the image is saved 
      in your current folder (default is None).
  
  max_image : boolean, optional
      Whether or not the image will be maximized within the space allocated 
      (default is True).

  textbox_filled_space : float (between 0 to 1), optional
      The portion of the slide width to be allocated to the textbox (default
      is 0.5).
  
  title_font_size : int, optional
      The font size of title text (default is 30).
  
  textbox_font_size : int, optional
      The font size of textbox text (default is 16).
  """
  
  slide = presentation_name.slides.add_slide(presentation_name.slide_layouts[5])

  buffer = Inches(0.25)
  
  if image_filepath == None:
    textbox_filled_space = 1
    buffer = 0
  
  if slide_text == "":
    textbox_filled_space = 0
    buffer = 0

  slide.shapes[0].text_frame.text = slide_title
  slide.shapes[0].text_frame.paragraphs[0].font.size = Pt(title_font_size)
  slide.shapes[0].text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

  left = slide.shapes[0].left
  top = slide.shapes[0].top + slide.shapes[0].height + Inches(0.25)
  width = round((slide.shapes[0].width - buffer) * textbox_filled_space)
  height = presentation_name.slide_height - top - Inches(0.25)

  img_left = left + width + buffer
  img_top = top
  img_width = round((slide.shapes[0].width - buffer) * (1 - textbox_filled_space))
  img_height = height
  
  textbox = slide.shapes.add_textbox(left,top,width,height)
  textbox.text = slide_text
  for p in range(len(textbox.text_frame.paragraphs)):
    makeParaBulletPointed(textbox.text_frame.paragraphs[p])
    textbox.text_frame.paragraphs[p].font.size = Pt(textbox_font_size)
  textbox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

  if image_filepath != None:
    slide.shapes.add_picture(image_filepath, left = img_left, top = img_top)
    if max_image == True:
      if slide.shapes[2].width - img_width > slide.shapes[2].height - img_height:
        ratio = img_width/slide.shapes[2].width
      else:
        ratio = img_height/slide.shapes[2].height
      slide.shapes[2].width = round(slide.shapes[2].width * ratio)
      slide.shapes[2].height = round(slide.shapes[2].height * ratio)
    slide.shapes[2].left = img_left + round((img_width - slide.shapes[2].width)/2)
    slide.shapes[2].top = img_top + round((img_height - slide.shapes[2].height)/2)

def save_deck(presentation_name, filepath="./Consultify.pptx"):
  """
  Saves the Presentation.

  Parameters
  ----------
  presentation_name : Presentation object
      The Presentation object created using the make_deck() function.

  filepath : str, optional
      The filepath where the resulting PowerPoint slide will be saved.
      Include the PowerPoint filename (default is "./Consultify.pptx").
  """
  if not filepath.endswith(".pptx"):
    filepath = filepath + ".pptx"
  presentation_name.save(filepath)